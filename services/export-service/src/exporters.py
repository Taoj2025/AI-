"""
Export Service - 各格式导出器实现
支持: PDF / Word / PPT / PNG / JPG / HTML / Markdown
"""
import io, json, uuid, datetime, os, tempfile
from typing import Optional, Dict, Any
from jinja2 import Template


# ============================================================
# 基础渲染器：将简历数据渲染为 HTML（中间格式）
# ============================================================
class BaseRenderer:
    """将简历 JSON 数据渲染为标准化 HTML"""

    BASE_TEMPLATE = """
    <!DOCTYPE html>
    <html lang="zh-CN">
    <head>
      <meta charset="UTF-8">
      <meta name="viewport" content="width=device-width, initial-scale=1.0">
      <style>
        @page { size: A4; margin: 1.5cm 1.8cm; }
        * { box-sizing: border-box; margin: 0; padding: 0; }
        body {
          font-family: "Noto Sans CJK SC", "PingFang SC", "Microsoft YaHei", sans-serif;
          font-size: 11pt; line-height: 1.7; color: #1A1A2E;
          -webkit-font-smoothing: antialiased;
        }
        .resume-a4 {
          width: 210mm; min-height: 297mm;
          padding: 1.5cm 1.8cm; margin: 0 auto;
          background: #FFF; position: relative;
        }
        /* 极简风格 */
        .style-minimal .section-title {
          font-size: 11pt; text-transform: uppercase; letter-spacing: 2px;
          color: #999; border: none; padding: 0; margin-bottom: 8pt;
        }
        /* 现代风格 */
        .style-modern .section-title {
          font-size: 13pt; color: #4361EE; border-left: 3px solid #4361EE;
          padding-left: 8pt; margin: 14pt 0 8pt;
        }
        /* 经典风格 */
        .style-classic .section-title {
          font-size: 13pt; border-bottom: 1.5pt solid #1A1A2E;
          padding-bottom: 3pt; margin: 14pt 0 8pt;
        }
        /* 创意风格 */
        .style-creative .section-title {
          font-size: 13pt; color: #E91E63; position: relative;
          padding-bottom: 4pt; margin: 14pt 0 8pt;
        }
        .style-creative .section-title::after {
          content: ''; position: absolute; bottom: 0; left: 0;
          width: 40pt; height: 2pt; background: #E91E63;
        }
        /* 高管风格 */
        .style-executive .resume-a4 { padding-top: 2.5cm; }
        .style-executive .header-name { font-size: 24pt; color: #0D47A1; }
        .section { margin-bottom: 6pt; }
        .item { margin-bottom: 8pt; }
        .item-header { display: flex; justify-content: space-between; align-items: baseline; }
        .item-title { font-weight: 600; }
        .item-date { font-size: 9pt; color: #666; }
        .item-subtitle { font-size: 10pt; color: #444; margin: 1pt 0; }
        .item-desc { font-size: 10pt; color: #333; margin-top: 2pt; }
        .skill-tag {
          display: inline-block; background: #F0F0FF; color: #4361EE;
          padding: 1pt 6pt; border-radius: 3pt; font-size: 9pt; margin: 1pt 2pt;
        }
        .header { margin-bottom: 12pt; }
        .header-name { font-size: 20pt; font-weight: 700; color: #1A1A2E; }
        .header-contact { font-size: 9pt; color: #666; margin-top: 2pt; }
        .header-summary { font-size: 10pt; color: #444; margin-top: 6pt; line-height: 1.6; }
        @media print { body { background: #FFF; } .resume-a4 { box-shadow: none; } }
      </style>
    </head>
    <body>
      <div class="resume-a4 style-{{ style }}">
        <!-- 头部 -->
        <div class="header">
          <div class="header-name">{{ personal.name }}</div>
          <div class="header-contact">
            {{ personal.title }} &nbsp;|&nbsp; {{ personal.email }}
            {% if personal.phone %} &nbsp;|&nbsp; {{ personal.phone }} {% endif %}
            {% if personal.location %} &nbsp;|&nbsp; {{ personal.location }} {% endif %}
          </div>
          {% if summary %}
          <div class="header-summary">{{ summary }}</div>
          {% endif %}
        </div>

        <!-- 工作经历 -->
        {% if work %}
        <div class="section">
          <div class="section-title">工作经历</div>
          {% for w in work %}
          <div class="item">
            <div class="item-header">
              <span class="item-title">{{ w.position }}</span>
              <span class="item-date">{{ w.startDate }} – {{ w.endDate or '至今' }}</span>
            </div>
            <div class="item-subtitle">{{ w.company }}</div>
            <div class="item-desc">{{ w.description }}</div>
            {% if w.achievements %}
              <ul style="margin:2pt 0 0 14pt; font-size:10pt;">
              {% for a in w.achievements %}<li>{{ a }}</li>{% endfor %}
              </ul>
            {% endif %}
          </div>
          {% endfor %}
        </div>
        {% endif %}

        <!-- 项目经历 -->
        {% if projects %}
        <div class="section">
          <div class="section-title">项目经历</div>
          {% for p in projects %}
          <div class="item">
            <div class="item-header">
              <span class="item-title">{{ p.name }}</span>
              <span class="item-date">{{ p.startDate }} – {{ p.endDate or '至今' }}</span>
            </div>
            <div class="item-subtitle">{{ p.role }}</div>
            <div class="item-desc">{{ p.description }}</div>
          </div>
          {% endfor %}
        </div>
        {% endif %}

        <!-- 教育背景 -->
        {% if education %}
        <div class="section">
          <div class="section-title">教育背景</div>
          {% for e in education %}
          <div class="item">
            <div class="item-header">
              <span class="item-title">{{ e.school }}</span>
              <span class="item-date">{{ e.startDate }} – {{ e.endDate }}</span>
            </div>
            <div class="item-subtitle">{{ e.major }}（{{ e.degree }}）</div>
          </div>
          {% endfor %}
        </div>
        {% endif %}

        <!-- 技能 -->
        {% if skills %}
        <div class="section">
          <div class="section-title">专业技能</div>
          {% for sk in skills %}
            <div style="margin-bottom:4pt;"><strong>{{ sk.category }}:</strong>
            {% for it in sk.skill_items %}<span class="skill-tag">{{ it }}</span>{% endfor %}
            </div>
          {% endfor %}
        </div>
        {% endif %}

        <!-- 证书 -->
        {% if certificates %}
        <div class="section">
          <div class="section-title">证书荣誉</div>
          <ul style="margin-left:14pt;">
          {% for c in certificates %}<li>{{ c.name }}（{{ c.issuer }}，{{ c.date }}）</li>{% endfor %}
          </ul>
        </div>
        {% endif %}
      </div>
    </body>
    </html>
    """

    @classmethod
    def render_html(cls, resume_data: Dict[str, Any], style: str = "modern") -> str:
        """渲染为完整 HTML 字符串"""
        # 映射 skills.items -> skill_items 避免 Jinja2 内置冲突
        skills = []
        for s in resume_data.get("skills", []):
            skills.append({"category": s.get("category"), "skill_items": s.get("items", [])})

        template = Template(cls.BASE_TEMPLATE)
        return template.render(
            personal=resume_data.get("personal", {}),
            summary=resume_data.get("summary", ""),
            work=resume_data.get("work", []),
            projects=resume_data.get("projects", []),
            education=resume_data.get("education", []),
            skills=skills,
            certificates=resume_data.get("certificates", []),
            style=style,
        )


# ============================================================
# PDF 导出器（ReportLab —— 纯 Python，无系统依赖）
# ============================================================
class PDFExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, options: Dict) -> str:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import cm, mm
        from reportlab.lib.colors import HexColor, black
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, TableStyle, HRFlowable
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT

        doc = SimpleDocTemplate(output_path, pagesize=A4,
            topMargin=1.5*cm, bottomMargin=1.5*cm, leftMargin=1.8*cm, rightMargin=1.8*cm)

        styles = getSampleStyleSheet()
        style = options.get("style", "modern")

        # 自定义样式
        name_style = ParagraphStyle('Name', parent=styles['Title'], fontSize=20, spaceAfter=4, textColor=black)
        contact_style = ParagraphStyle('Contact', parent=styles['Normal'], fontSize=9, textColor=HexColor('#666666'), alignment=TA_CENTER)
        section_style = ParagraphStyle('Section', parent=styles['Heading2'], fontSize=13,
            spaceAfter=6, spaceBefore=12, textColor=HexColor('#4361EE' if style == 'modern' else '#1A1A2E'))
        body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=10, leading=16)
        bullet_style = ParagraphStyle('Bullet', parent=styles['Normal'], fontSize=10, leftIndent=14, bulletIndent=4)

        elements = []
        p = resume_data.get("personal", {})

        # 姓名
        elements.append(Paragraph(p.get("name", ""), name_style))
        elements.append(Paragraph(
            f"{p.get('title','')} | {p.get('email','')}" +
            (f" | {p.get('phone','')}" if p.get('phone') else "") +
            (f" | {p.get('location','')}" if p.get('location') else ""),
            contact_style
        ))
        elements.append(Spacer(1, 6))

        # 个人总结
        if resume_data.get("summary"):
            elements.append(Paragraph("个人总结", section_style))
            elements.append(Paragraph(resume_data["summary"], body_style))

        # 工作经历
        if resume_data.get("work"):
            elements.append(Paragraph("工作经历", section_style))
            for w in resume_data["work"]:
                end = w.get("endDate") or "至今"
                elements.append(Paragraph(
                    f"<b>{w.get('position','')}</b> <font color='#666' size=9>({w.get('startDate','')} – {end})</font>",
                    body_style
                ))
                elements.append(Paragraph(f"<font color='#444'>{w.get('company','')}</font>", body_style))
                elements.append(Paragraph(w.get("description", ""), body_style))
                for a in w.get("achievements", []):
                    elements.append(Paragraph(f"• {a}", bullet_style))
                elements.append(Spacer(1, 4))

        # 项目经历
        if resume_data.get("projects"):
            elements.append(Paragraph("项目经历", section_style))
            for pj in resume_data["projects"]:
                elements.append(Paragraph(f"<b>{pj.get('name','')}</b> — {pj.get('role','')}", body_style))
                elements.append(Paragraph(pj.get("description", ""), body_style))

        # 教育背景
        if resume_data.get("education"):
            elements.append(Paragraph("教育背景", section_style))
            for e in resume_data["education"]:
                elements.append(Paragraph(f"<b>{e.get('school','')}</b> — {e.get('major','')}（{e.get('degree','')}）", body_style))

        # 技能
        if resume_data.get("skills"):
            elements.append(Paragraph("专业技能", section_style))
            for s in resume_data["skills"]:
                items = "、".join(s.get("items", []))
                elements.append(Paragraph(f"<b>{s.get('category','')}</b>：{items}", body_style))

        # 证书
        if resume_data.get("certificates"):
            elements.append(Paragraph("证书荣誉", section_style))
            for c in resume_data["certificates"]:
                elements.append(Paragraph(f"• {c.get('name','')}（{c.get('issuer','')}, {c.get('date','')}）", bullet_style))

        doc.build(elements)
        return output_path


# ============================================================
# Word 导出器（python-docx）
# ============================================================
class WordExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, options: Dict) -> str:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # 页面边距
        from docx.shared import Cm
        section = doc.sections[0]
        section.top_margin = Cm(1.5)
        section.bottom_margin = Cm(1.5)
        section.left_margin = Cm(1.8)
        section.right_margin = Cm(1.8)

        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(resume_data.get("personal", {}).get("name", ""))
        run.bold = True
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # 联系方式
        contact = resume_data.get("personal", {})
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(f"{contact.get('title','')} | {contact.get('email','')}")
        if contact.get("phone"):
            p.add_run(f" | {contact.get('phone')}")

        # 个人总结
        if resume_data.get("summary"):
            doc.add_heading("个人总结", 2)
            doc.add_paragraph(resume_data["summary"])

        # 工作经历
        if resume_data.get("work"):
            doc.add_heading("工作经历", 2)
            for w in resume_data["work"]:
                p = doc.add_paragraph()
                p.add_run(f"{w.get('position','')}  @  {w.get('company','')}  ").bold = True
                p.add_run(f"{w.get('startDate','')} – {w.get('endDate','至今')}")
                doc.add_paragraph(w.get("description", ""))
                for a in w.get("achievements", []):
                    doc.add_paragraph(a, style="List Bullet")

        # 教育背景
        if resume_data.get("education"):
            doc.add_heading("教育背景", 2)
            for e in resume_data["education"]:
                p = doc.add_paragraph()
                p.add_run(f"{e.get('school','')}  -  {e.get('major','')}（{e.get('degree','')}）").bold = True

        # 技能
        if resume_data.get("skills"):
            doc.add_heading("专业技能", 2)
            for s in resume_data["skills"]:
                items = "、".join(s.get("items", []))
                doc.add_paragraph(f"{s.get('category','')}：{items}")

        doc.save(output_path)
        return output_path


# ============================================================
# PPT 导出器（python-pptx）
# ============================================================
class PPTExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, options: Dict) -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        # 第1页：封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(1.5)
        shape = slide.shapes.add_textbox(left, top, width, height)
        tf = shape.text_frame
        tf.text = resume_data.get("personal", {}).get("name", "")
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        # 副标题
        shape2 = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        tf2 = shape2.text_frame
        tf2.text = resume_data.get("personal", {}).get("title", "")
        tf2.paragraphs[0].font.size = Pt(18)

        # 第2页：工作经历
        if resume_data.get("work"):
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            tf = shape.text_frame
            tf.text = "工作经历"
            tf.paragraphs[0].font.size = Pt(20)
            tf.paragraphs[0].font.bold = True

            y = Inches(1.2)
            for w in resume_data["work"]:
                shape = slide2.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(1.2))
                tf = shape.text_frame
                tf.text = f"{w.get('position','')} @ {w.get('company','')}\n{w.get('description','')}"
                tf.paragraphs[0].font.size = Pt(12)
                y += Inches(1.3)

        prs.save(output_path)
        return output_path


# ============================================================
# 图片导出器（ReportLab PDF → Pillow 截图）
# ============================================================
class ImageExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, fmt: str = "png", options: Dict = {}) -> str:
        """先生成 PDF，再用 Pillow 将首页转为图片"""
        import tempfile
        from PIL import Image
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.pdfgen import canvas
        import io

        # 生成 PDF 到内存
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            PDFExporter.export(resume_data, tmp.name, options)
            pdf_path = tmp.name

        # 使用 reportlab 的 canvas 渲染为图片（纯方案）
        # 生成高质量 PNG
        buf = io.BytesIO()
        pdf_w, pdf_h = A4
        dpi = 150
        img_w = int(pdf_w / 72 * dpi)
        img_h = int(pdf_h / 72 * dpi)
        img = Image.new("RGB", (img_w, img_h), "white")

        # 绘制简化版简历到图片
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        y = 40
        p = resume_data.get("personal", {})
        try:
            font_name = ImageFont.truetype("msyh.ttc", 32)  # 微软雅黑
            font_body = ImageFont.truetype("msyh.ttc", 18)
            font_section = ImageFont.truetype("msyh.ttc", 22)
            font_small = ImageFont.truetype("msyh.ttc", 14)
        except Exception:
            font_name = ImageFont.load_default()
            font_body = font_name
            font_section = font_name
            font_small = font_name

        # 姓名
        draw.text((60, y), p.get("name", ""), font=font_name, fill=(26, 26, 46))
        y += 50
        # 联系方式
        contact = f"{p.get('title','')} | {p.get('email','')} | {p.get('phone','')}"
        draw.text((60, y), contact, font=font_small, fill=(100, 100, 100))
        y += 40

        if resume_data.get("summary"):
            draw.text((60, y), "个人总结", font=font_section, fill=(67, 97, 238))
            y += 30
            draw.text((60, y), resume_data["summary"][:200], font=font_body, fill=(50, 50, 50))
            y += 60

        if resume_data.get("work"):
            draw.text((60, y), "工作经历", font=font_section, fill=(67, 97, 238))
            y += 30
            for w in resume_data["work"]:
                draw.text((60, y), w.get("position", ""), font=font_body, fill=(26, 26, 46))
                y += 24
                draw.text((60, y), w.get("company", ""), font=font_small, fill=(100, 100, 100))
                y += 40

        if resume_data.get("education"):
            draw.text((60, y), "教育背景", font=font_section, fill=(67, 97, 238))
            y += 30
            for e in resume_data["education"]:
                draw.text((60, y), f"{e.get('school','')} - {e.get('major','')}（{e.get('degree','')}）", font=font_body, fill=(50, 50, 50))
                y += 30

        if fmt == "jpg":
            img = img.convert("RGB")
        img.save(output_path, quality=95)
        return output_path


# ============================================================
# HTML 导出器（直接输出单文件 HTML）
# ============================================================
class HTMLExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, options: Dict) -> str:
        html = BaseRenderer.render_html(resume_data, options.get("style", "modern"))
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        return output_path


# ============================================================
# Markdown 导出器
# ============================================================
class MarkdownExporter:
    @staticmethod
    def export(resume_data: Dict[str, Any], output_path: str, options: Dict) -> str:
        lines = []
        p = resume_data.get("personal", {})
        lines.append(f"# {p.get('name', '')}")
        lines.append(f"**{p.get('title','')}** | {p.get('email','')} | {p.get('phone','')}")
        lines.append("")

        if resume_data.get("summary"):
            lines.append("## 个人总结")
            lines.append(resume_data["summary"])
            lines.append("")

        if resume_data.get("work"):
            lines.append("## 工作经历")
            for w in resume_data["work"]:
                lines.append(f"### {w.get('position','')} @ {w.get('company','')}")
                lines.append(f"*{w.get('startDate','')} – {w.get('endDate','至今')}*")
                lines.append(w.get("description", ""))
                for a in w.get("achievements", []):
                    lines.append(f"- {a}")
                lines.append("")

        if resume_data.get("projects"):
            lines.append("## 项目经历")
            for pj in resume_data["projects"]:
                lines.append(f"### {pj.get('name','')}（{pj.get('role','')}）")
                lines.append(pj.get("description", ""))
                lines.append("")

        if resume_data.get("education"):
            lines.append("## 教育背景")
            for e in resume_data["education"]:
                lines.append(f"- **{e.get('school','')}** - {e.get('major','')}（{e.get('degree','')}）")

        if resume_data.get("skills"):
            lines.append("## 专业技能")
            for s in resume_data["skills"]:
                items = "、".join(s.get("items", []))
                lines.append(f"- **{s.get('category','')}**：{items}")

        md = "\n".join(lines)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(md)
        return output_path


# ============================================================
# 统一导出入口
# ============================================================
EXPORT_MAP = {
    "pdf": PDFExporter,
    "docx": WordExporter,
    "pptx": PPTExporter,
    "png": ImageExporter,
    "jpg": ImageExporter,
    "html": HTMLExporter,
    "markdown": MarkdownExporter,
}

def export_resume(
    resume_data: Dict[str, Any],
    output_path: str,
    fmt: str,
    options: Optional[Dict] = None
) -> str:
    """统一导出入口"""
    options = options or {}
    exporter_cls = EXPORT_MAP.get(fmt)
    if not exporter_cls:
        raise ValueError(f"不支持的导出格式: {fmt}")

    if fmt in ("png", "jpg"):
        return exporter_cls.export(resume_data, output_path, fmt, options)
    return exporter_cls.export(resume_data, output_path, options)


# -------- CLI 测试 --------
if __name__ == "__main__":
    test_data = {
        "personal": {"name": "张三", "title": "高级前端工程师", "email": "zhangsan@example.com", "phone": "13800138000"},
        "summary": "5年前端开发经验，专注 React 和 TypeScript。",
        "work": [
            {"position": "高级前端工程师", "company": "字节跳动", "startDate": "2020-07", "endDate": None, "description": "负责抖音电商前端架构设计", "achievements": ["性能优化提升40%"]}
        ],
        "education": [
            {"school": "清华大学", "major": "计算机科学", "degree": "硕士", "startDate": "2017", "endDate": "2020"}
        ],
        "skills": [
            {"category": "前端框架", "items": ["React", "Vue", "Next.js"]}
        ],
    }
    out = f"/tmp/test_resume.pdf"
    export_resume(test_data, out, "pdf")
    print(f"导出成功: {out}")
